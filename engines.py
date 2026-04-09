import os
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
import google.generativeai as genai
import chromadb
import time
from langchain_community.document_loaders import PyPDFLoader
from langchain_text_splitters import RecursiveCharacterTextSplitter
from pptx import Presentation
from pptx.util import Inches

class KnowledgeEngine:
    def __init__(self):
        self.chroma_client = chromadb.Client()
        try:
            self.collection = self.chroma_client.create_collection(name="strat_vault")
        except:
            self.collection = self.chroma_client.get_collection(name="strat_vault")
        self.text_splitter = RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=100)

    def ingest_pdf(self, file_object):
        try:
            with open("temp.pdf", "wb") as f:
                f.write(file_object.getbuffer())
            loader = PyPDFLoader("temp.pdf")
            pages = loader.load()
            chunks = self.text_splitter.split_documents(pages)
            for i, chunk in enumerate(chunks):
                self.collection.add(documents=[chunk.page_content], ids=[f"id_{os.urandom(4).hex()}"])
            os.remove("temp.pdf")
            return len(chunks)
        except Exception as e:
            return f"Error indexing PDF: {str(e)}"

    def query(self, text):
        try:
            results = self.collection.query(query_texts=[text], n_results=3)
            return " ".join(results['documents'][0]) if results['documents'] else ""
        except:
            return ""

class DataEngine:
    def clean_and_load(self, file_object):
        try:
            if file_object.name.endswith('.csv'):
                df = pd.read_csv(file_object)
            else:
                df = pd.read_excel(file_object)
            
            df.columns = df.columns.astype(str).str.strip()
            df = df.dropna(how='all', axis=0).dropna(how='all', axis=1)
            
            for col in df.columns:
                if df[col].dtype == 'object':
                    try:
                        df[col] = df[col].replace(r'[\$,%]', '', regex=True).astype(float)
                    except:
                        pass 
            return df
        except Exception as e:
            return None

    def run_ai_realise_audit(self, df):
        """Lancia AI Realise: Data Maturity Scoring"""
        if df is None or df.empty:
            return {"score": 0, "status": "NO DATA", "color": "red"}
        
        completeness = (1 - df.isnull().mean().mean()) * 100
        numeric_density = (len(df.select_dtypes(include=[np.number]).columns) / len(df.columns)) * 100
        
        maturity_score = round((completeness + numeric_density) / 2, 1)
        
        if maturity_score > 85:
            return {"score": maturity_score, "status": "STRATEGICALLY READY", "color": "green"}
        elif maturity_score > 60:
            return {"score": maturity_score, "status": "DATA REMEDIATION ADVISED", "color": "orange"}
        else:
            return {"score": maturity_score, "status": "CRITICAL DATA GAPS", "color": "red"}

    def analyze_and_plot(self, df):
        numeric = df.select_dtypes(include=[np.number])
        if numeric.empty:
            return {"Status": "No numeric trends found."}, []
        
        stats, charts = {}, []
        sns.set_theme(style="whitegrid")
        
        for col in numeric.columns[:3]:
            vals = numeric[col].dropna()
            if len(vals) > 1:
                trend = np.polyfit(range(len(vals)), vals, 1)[0]
                stats[col] = {"mean": round(vals.mean(), 2), "trend": round(trend, 4)}
                
                plt.figure(figsize=(6, 4))
                sns.lineplot(data=vals, color='#074050', linewidth=2, marker='o')
                plt.title(f"Client Data Analysis: {col}")
                path = f"{col.replace(' ', '_')}_viz.png"
                plt.savefig(path, bbox_inches='tight')
                plt.close()
                charts.append(path)
        return stats, charts

    def generate_sensitivity(self, target_gain):
        x = np.linspace(0, target_gain * 1.5, 20)
        y = 15 + (x * 0.82) 
        
        plt.figure(figsize=(8, 4))
        plt.plot(x, y, color='#074050', linewidth=3, label="Profit Trajectory")
        plt.axvline(target_gain, color='#d9534f', linestyle='--', label=f"Target: {target_gain}%")
        plt.fill_between(x, y-2, y+2, alpha=0.1, color='#074050')
        plt.title("Strategic Sensitivity: Efficiency vs. Projected Margin", fontsize=10)
        plt.xlabel("OpEx Efficiency Gain (%)")
        plt.ylabel("Projected Net Margin (%)")
        plt.legend()
        plt.grid(alpha=0.3)
        
        path = "sensitivity_plot.png"
        plt.savefig(path, bbox_inches='tight')
        plt.close()
        return path

class StratOS_Orchestrator:
    def __init__(self, kb, api_key):
        genai.configure(api_key=api_key)
        self.kb = kb
        self.model_stack = ['gemini-2.0-flash', 'models/gemini-2.0-flash', 'gemini-1.5-flash']

    def get_persona_database(self):
        """Core Lancia Service Lines"""
        return {
            "Retail": "Consumer Analyst. Focus: Inventory churn, CAC/LTV, and omnichannel conversion.",
            "Logistics": "Supply Chain Architect. Focus: Lead times, fuel volatility, and port congestion.",
            "Utilities": "Infrastructure Strategist. Focus: Grid stability, ESG, and CAPEX modeling.",
            "Manufacturing": "Lean Specialist. Focus: OEE, waste reduction, and unit-cost optimization."
        }

    def run_debate(self, problem, stats, industry):
        lancia_context = self.get_persona_database().get(industry, "Senior Strategy Partner")
        agents = {
            "The Visionary (Growth)": f"Focus on market capture. Industry context: {lancia_context}",
            "The Conservator (Risk)": f"Focus on cash flow and compliance. Industry context: {lancia_context}",
            "The Operationalist (Execution)": f"Focus on scalability. Industry context: {lancia_context}"
        }
        
        transcript = ""
        context = self.kb.query(problem)
        
        for name, persona in agents.items():
            prompt = (f"ROLE: {name}. {persona}\nINDUSTRY: {industry}\n"
                      f"DATA: {stats}\nCONTEXT: {context}\nTASK: 2-sentence solution for: {problem}")
            response = self.call_gemini(prompt)
            transcript += f"**{name}**: {response}\n\n"
        
        synthesis = (f"Synthesize this debate into a unified 3-pillar strategy for {industry}:\n\n{transcript}\n\n"
                     "Deliverable: Pyramid Principle (Governing Thought + 3 MECE Pillars).")
        return transcript, self.call_gemini(synthesis)

    def generate_lancia_roadmap(self, strategy_text):
        """Lancia 90-Day Implementation Methodology"""
        prompt = (f"Act as a Lancia Transformation Lead. Based on: {strategy_text}\n"
                  "Create a 90-Day Roadmap table: Month 1 (Mobilize), Month 2 (Execute), Month 3 (Realize). "
                  "Provide 2 clear KPIs per month.")
        return self.call_gemini(prompt)

    def run_red_team_audit(self, strategy):
        prompt = f"Act as a cynical Auditor. Find 3 structural failure points in: {strategy}. Provide mitigations."
        return self.call_gemini(prompt)

    def call_gemini(self, prompt):
        for model_name in self.model_stack:
            try:
                model = genai.GenerativeModel(model_name)
                time.sleep(1) 
                response = model.generate_content(prompt)
                if response and response.text:
                    return response.text
            except:
                continue
        return "CRITICAL ERROR: Strategy Engine unavailable."

class ReportEngine:
    def run_results365_check(self, strategy_output):
        """Lancia Results365: Performance Health Diagnostic"""
        return (f"Act as the Results365 Health Auditor. Evaluate this strategy: {strategy_output}\n"
                "Identify 3 risks that would cause project slippage. Rate delivery confidence (0-100%).")

    def create_deck(self, summary, charts, roadmap=""):
        prs = Presentation()
        # Slide 1: Synthesis
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Executive Strategy Mandate"
        slide.shapes.placeholders[1].text = summary[:1500]
        # Slide 2: Roadmap
        if roadmap:
            slide_r = prs.slides.add_slide(prs.slide_layouts[1])
            slide_r.shapes.title.text = "90-Day Implementation Roadmap"
            slide_r.shapes.placeholders[1].text = roadmap[:1500]
        # Chart Slides
        for img in charts:
            if os.path.exists(img):
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(img, Inches(0.5), Inches(1.5), width=Inches(9))
        
        path = "Lancia_Strategy_Report.pptx"
        prs.save(path)
        return path
